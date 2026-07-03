import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

class PresentationBuilder:
    def __init__(self, theme='navy'):
        """
        Initializes the PresentationBuilder with a selected theme.
        Themes:
            - 'navy': Deep Slate/Navy background for title, light for content, amber accent.
            - 'forest': Deep Forest Green background for title, light for content, emerald accent.
            - 'slate': Slate Grey background for title, light for content, rose/coral accent.
        """
        self.prs = Presentation()
        # Set 16:9 widescreen dimensions
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
        # Configure theme color schemes
        self.themes = {
            'navy': {
                'dark_bg': RGBColor(15, 23, 42),       # Slate 900
                'light_bg': RGBColor(255, 255, 255),   # White
                'title_text': RGBColor(15, 23, 42),    # Slate 900 for light bg
                'body_text': RGBColor(51, 65, 85),     # Slate 700
                'accent': RGBColor(245, 158, 11),      # Amber 500
                'accent_light': RGBColor(254, 243, 199) # Amber 100
            },
            'forest': {
                'dark_bg': RGBColor(6, 78, 59),        # Green 900
                'light_bg': RGBColor(255, 255, 255),   # White
                'title_text': RGBColor(6, 78, 59),     # Green 900
                'body_text': RGBColor(20, 83, 45),     # Green 800
                'accent': RGBColor(16, 185, 129),      # Emerald 500
                'accent_light': RGBColor(209, 250, 229) # Emerald 100
            },
            'slate': {
                'dark_bg': RGBColor(30, 41, 59),       # Slate 800
                'light_bg': RGBColor(250, 250, 250),   # Off-white
                'title_text': RGBColor(15, 23, 42),    # Slate 900
                'body_text': RGBColor(71, 85, 105),    # Slate 600
                'accent': RGBColor(244, 63, 94),       # Rose 500
                'accent_light': RGBColor(FFE4E6) if False else RGBColor(255, 228, 230) # Rose 100
            }
        }
        
        self.theme_name = theme if theme in self.themes else 'navy'
        self.colors = self.themes[self.theme_name]
        self.font_name = 'Microsoft JhengHei' # Default professional Chinese/Bilingual font

    def _set_slide_background(self, slide, color):
        """Sets a solid color background for the slide."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_title_slide(self, title, subtitle):
        """Adds a premium dark-themed title slide."""
        # Use blank layout (usually layout index 6 in default templates, but let's use layout[6] or layout[5])
        # To be safe and independent of template layout order, we can use the blank layout (usually index 6)
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # Set dark background
        self._set_slide_background(slide, self.colors['dark_bg'])
        
        # Single Text box for both title and subtitle to prevent overlapping
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(4.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        
        # Title paragraph
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = self.font_name
        p_title.font.size = Pt(44)
        p_title.font.bold = True
        p_title.font.color.rgb = RGBColor(255, 255, 255)
        p_title.space_after = Pt(20)
        
        # Subtitle paragraph
        if subtitle:
            p_sub = tf.add_paragraph()
            p_sub.text = subtitle
            p_sub.font.name = self.font_name
            p_sub.font.size = Pt(22)
            p_sub.font.color.rgb = self.colors['accent']
            
        return slide

    def add_section_divider(self, title, description=None):
        """Adds a clean transitional section divider slide."""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # Set dark background
        self._set_slide_background(slide, self.colors['dark_bg'])
        
        # Add decorative top accent line
        # Left, Top, Width, Height
        line = slide.shapes.add_shape(
            1, # MSO_SHAPE.RECTANGLE = 1
            Inches(1.0), Inches(1.8), Inches(2.0), Inches(0.08)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.colors['accent']
        line.line.color.rgb = self.colors['accent']
        
        # Title text box
        box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(4.0))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = self.font_name
        p_title.font.size = Pt(36)
        p_title.font.bold = True
        p_title.font.color.rgb = RGBColor(255, 255, 255)
        p_title.space_after = Pt(16)
        
        if description:
            p_desc = tf.add_paragraph()
            p_desc.text = description
            p_desc.font.name = self.font_name
            p_desc.font.size = Pt(18)
            p_desc.font.color.rgb = RGBColor(203, 213, 225) # Slate 300
            
        return slide

    def add_content_slide(self, title, points):
        """Adds a standard content slide with a title and bullet points."""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # Set light background
        self._set_slide_background(slide, self.colors['light_bg'])
        
        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.0))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_top = Inches(0)
        tf_title.margin_bottom = Inches(0)
        tf_title.margin_left = Inches(0)
        tf_title.margin_right = Inches(0)
        
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.name = self.font_name
        p_title.font.size = Pt(32)
        p_title.font.bold = True
        p_title.font.color.rgb = self.colors['title_text']
        
        # Content Box
        content_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(4.5))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True
        tf_content.margin_top = Inches(0)
        tf_content.margin_bottom = Inches(0)
        tf_content.margin_left = Inches(0)
        tf_content.margin_right = Inches(0)
        
        for i, pt in enumerate(points):
            p = tf_content.paragraphs[0] if i == 0 else tf_content.add_paragraph()
            p.text = "•  " + pt
            p.font.name = self.font_name
            p.font.size = Pt(18)
            p.font.color.rgb = self.colors['body_text']
            p.space_after = Pt(16)
            
        return slide

    def add_split_slide(self, title, left_title, left_points, right_title, right_points):
        """Adds a two-column comparison slide."""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # Set light background
        self._set_slide_background(slide, self.colors['light_bg'])
        
        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.0))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_top = Inches(0)
        tf_title.margin_bottom = Inches(0)
        tf_title.margin_left = Inches(0)
        tf_title.margin_right = Inches(0)
        
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.name = self.font_name
        p_title.font.size = Pt(32)
        p_title.font.bold = True
        p_title.font.color.rgb = self.colors['title_text']
        
        # Left Column Box
        left_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.2), Inches(4.5))
        tf_left = left_box.text_frame
        tf_left.word_wrap = True
        tf_left.margin_top = Inches(0)
        tf_left.margin_bottom = Inches(0)
        tf_left.margin_left = Inches(0)
        tf_left.margin_right = Inches(0)
        
        # Left Column Title
        p_lt = tf_left.paragraphs[0]
        p_lt.text = left_title
        p_lt.font.name = self.font_name
        p_lt.font.size = Pt(22)
        p_lt.font.bold = True
        p_lt.font.color.rgb = self.colors['title_text']
        p_lt.space_after = Pt(14)
        
        # Left Column Bullet Points
        for pt in left_points:
            p = tf_left.add_paragraph()
            p.text = "•  " + pt
            p.font.name = self.font_name
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['body_text']
            p.space_after = Pt(12)
            
        # Right Column Box
        right_box = slide.shapes.add_textbox(Inches(7.133), Inches(2.2), Inches(5.2), Inches(4.5))
        tf_right = right_box.text_frame
        tf_right.word_wrap = True
        tf_right.margin_top = Inches(0)
        tf_right.margin_bottom = Inches(0)
        tf_right.margin_left = Inches(0)
        tf_right.margin_right = Inches(0)
        
        # Right Column Title
        p_rt = tf_right.paragraphs[0]
        p_rt.text = right_title
        p_rt.font.name = self.font_name
        p_rt.font.size = Pt(22)
        p_rt.font.bold = True
        p_rt.font.color.rgb = self.colors['accent']
        p_rt.space_after = Pt(14)
        
        # Right Column Bullet Points
        for pt in right_points:
            p = tf_right.add_paragraph()
            p.text = "•  " + pt
            p.font.name = self.font_name
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['body_text']
            p.space_after = Pt(12)
            
        return slide

    def add_metric_slide(self, title, metric_number, metric_label, description_points):
        """Adds a slide focused on displaying a high-impact metric on the left and description on the right."""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # Set light background
        self._set_slide_background(slide, self.colors['light_bg'])
        
        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.0))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_top = Inches(0)
        tf_title.margin_bottom = Inches(0)
        tf_title.margin_left = Inches(0)
        tf_title.margin_right = Inches(0)
        
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.name = self.font_name
        p_title.font.size = Pt(32)
        p_title.font.bold = True
        p_title.font.color.rgb = self.colors['title_text']
        
        # Left Column: Big Metric Box
        metric_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(5.2), Inches(4.0))
        tf_metric = metric_box.text_frame
        tf_metric.word_wrap = True
        tf_metric.margin_top = Inches(0)
        tf_metric.margin_bottom = Inches(0)
        tf_metric.margin_left = Inches(0)
        tf_metric.margin_right = Inches(0)
        
        # Big Number
        p_num = tf_metric.paragraphs[0]
        p_num.text = metric_number
        p_num.font.name = self.font_name
        p_num.font.size = Pt(84)
        p_num.font.bold = True
        p_num.font.color.rgb = self.colors['accent']
        p_num.space_after = Pt(10)
        
        # Metric Label
        p_label = tf_metric.add_paragraph()
        p_label.text = metric_label
        p_label.font.name = self.font_name
        p_label.font.size = Pt(20)
        p_label.font.bold = True
        p_label.font.color.rgb = self.colors['title_text']
        
        # Right Column: Description
        desc_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.4), Inches(5.5), Inches(4.0))
        tf_desc = desc_box.text_frame
        tf_desc.word_wrap = True
        tf_desc.margin_top = Inches(0)
        tf_desc.margin_bottom = Inches(0)
        tf_desc.margin_left = Inches(0)
        tf_desc.margin_right = Inches(0)
        
        for i, pt in enumerate(description_points):
            p = tf_desc.paragraphs[0] if i == 0 else tf_desc.add_paragraph()
            p.text = "•  " + pt
            p.font.name = self.font_name
            p.font.size = Pt(18)
            p.font.color.rgb = self.colors['body_text']
            p.space_after = Pt(14)
            
        return slide

    def add_chart_slide(self, title, chart_title, categories, series_name, series_data, points):
        """Adds a slide with bullet points on the left and a native column chart on the right."""
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self._set_slide_background(slide, self.colors['light_bg'])
        
        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.0))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_top = Inches(0)
        tf_title.margin_bottom = Inches(0)
        tf_title.margin_left = Inches(0)
        tf_title.margin_right = Inches(0)
        
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.name = self.font_name
        p_title.font.size = Pt(32)
        p_title.font.bold = True
        p_title.font.color.rgb = self.colors['title_text']
        
        # Left Column Box for Points
        content_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.2), Inches(4.5))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True
        tf_content.margin_top = Inches(0)
        tf_content.margin_bottom = Inches(0)
        tf_content.margin_left = Inches(0)
        tf_content.margin_right = Inches(0)
        
        for i, pt in enumerate(points):
            p = tf_content.paragraphs[0] if i == 0 else tf_content.add_paragraph()
            p.text = "•  " + pt
            p.font.name = self.font_name
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['body_text']
            p.space_after = Pt(12)
            
        # Right Column: Chart
        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series(series_name, series_data)
        
        x, y, cx, cy = Inches(6.8), Inches(2.2), Inches(5.5), Inches(4.5)
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        )
        chart = chart_shape.chart
        chart.has_legend = False
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_title
        chart.chart_title.text_frame.paragraphs[0].font.name = self.font_name
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = self.colors['title_text']
        
        # Style the series
        if len(chart.series) > 0:
            series = chart.series[0]
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = self.colors['accent']
            
        return slide

    def add_image_slide(self, title, image_path, points):
        """Adds a slide with bullet points on the left and an image on the right."""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self._set_slide_background(slide, self.colors['light_bg'])
        
        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.0))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_top = Inches(0)
        tf_title.margin_bottom = Inches(0)
        tf_title.margin_left = Inches(0)
        tf_title.margin_right = Inches(0)
        
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.name = self.font_name
        p_title.font.size = Pt(32)
        p_title.font.bold = True
        p_title.font.color.rgb = self.colors['title_text']
        
        # Left Column Box for Points
        content_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.2), Inches(4.5))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True
        tf_content.margin_top = Inches(0)
        tf_content.margin_bottom = Inches(0)
        tf_content.margin_left = Inches(0)
        tf_content.margin_right = Inches(0)
        
        for i, pt in enumerate(points):
            p = tf_content.paragraphs[0] if i == 0 else tf_content.add_paragraph()
            p.text = "•  " + pt
            p.font.name = self.font_name
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['body_text']
            p.space_after = Pt(12)
            
        # Right Column: Image
        try:
            slide.shapes.add_picture(image_path, Inches(6.8), Inches(2.2), width=Inches(5.5))
        except Exception as e:
            print(f"Warning: Could not add picture {image_path}: {e}")
            # Fallback text box if image fails
            fallback_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.2), Inches(5.5), Inches(4.5))
            fallback_box.text_frame.text = f"[Image: {os.path.basename(image_path)}]"
            
        return slide

    def save(self, filepath):
        """Saves the generated presentation to the specified filepath."""
        # Ensure target directory exists
        os.makedirs(os.path.dirname(os.path.abspath(filepath)), exist_ok=True)
        self.prs.save(filepath)
        print(f"Presentation saved successfully to: {filepath}")
